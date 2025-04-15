import json
import base64
import uuid
import logging
import boto3
import os
import pdfplumber
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger()

# === AWS Clients ===
s3 = boto3.client("s3", region_name="us-east-1")
bedrock = boto3.client("bedrock-runtime", region_name="us-east-1")
BUCKET_NAME = 'sowbucketcreation'  # ✅ Updated bucket name


def build_prompt(proposal_text):
    return f"""
You are a senior consultant and technical writer. 
You are tasked with generating a professional and contract-ready Statement of Work (SoW) for a 
technology project.

Use the following structure to create the SoW. Use professional, legal, and clear language 
appropriate for enterprise agreements.

Include these sections with rich detail:

1. DURATION
2. SERVICES AND DELIVERABLES
3. IMPLEMENTATION TIMELINE
4. ACCEPTANCE CRITERIA
5. GOVERNANCE AND MONITORING
6. TEAM & ROLES
7. CHARGES AND PAYMENT SCHEDULE
8. ASSUMPTIONS AND EXCLUSIONS
9. DATA PROTECTION (GDPR)
10. SIGN-OFF SECTION

Now generate a Statement of Work based on the following proposal:
{proposal_text}

Output the full SoW starting with the title: Statement of Work for [Project Title] and format each 
section clearly using titles and line breaks. Use legal and business terminology throughout.
"""


def generate_sow_from_prompt(prompt):
    try:
        response = bedrock.invoke_model(
            modelId="amazon.titan-text-express-v1",  # ✅ Titan Text G1 - Express for generation
            contentType="application/json",
            accept="application/json",
            body=json.dumps({
                "inputText": prompt,
                "textGenerationConfig": {
                    "maxTokenCount": 1000,
                    "temperature": 0.5,
                    "topP": 1
                }
            })
        )
        result = json.loads(response["body"].read())
        sow = result["results"][0]["outputText"]
        return sow.strip()
    except Exception as e:
        logger.error(f"Error from Bedrock: {e}")
        return f"⚠️ Bedrock error: {str(e)}"


def extract_pdf(file_bytes):
    with open("temp_local.pdf", "wb") as f:
        f.write(file_bytes)
    text = ""
    with pdfplumber.open("temp_local.pdf") as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text


def extract_pptx(file_bytes):
    with open("temp_local.pptx", "wb") as f:
        f.write(file_bytes)
    prs = Presentation("temp_local.pptx")
    return "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    )


def extract_text(file_path):
    _, ext = os.path.splitext(file_path)
    with open(file_path, "rb") as f:
        file_bytes = f.read()

    if ext.lower() == ".pdf":
        return extract_pdf(file_bytes)
    elif ext.lower() == ".pptx":
        return extract_pptx(file_bytes)
    else:
        return file_bytes.decode("utf-8")


def upload_to_s3(key, data_bytes):
    s3.put_object(Bucket=BUCKET_NAME, Key=key, Body=data_bytes)
    logger.info(f"Uploaded to s3://{BUCKET_NAME}/{key}")


def generate_presigned_url(key, expiration=3600):
    return s3.generate_presigned_url('get_object', Params={'Bucket': BUCKET_NAME, 'Key': key}, ExpiresIn=expiration)


def process_file_and_generate_sow(input_file_path: str):
    proposal_id = str(uuid.uuid4())
    proposal_text = extract_text(input_file_path)

    if not proposal_text.strip():
        raise ValueError("No valid proposal text extracted.")

    prompt = build_prompt(proposal_text)
    sow = generate_sow_from_prompt(prompt)

    # Save to S3
    upload_to_s3(f'sows/{proposal_id}.txt', sow.encode("utf-8"))
    upload_to_s3(f'proposals/{proposal_id}.txt', proposal_text.encode("utf-8"))

    # Generate download URL
    url = generate_presigned_url(f'sows/{proposal_id}.txt')
    return sow, url
